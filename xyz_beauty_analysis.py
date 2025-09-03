#!/usr/bin/env python
# coding: utf-8

# In[1]:


#!/usr/bin/env python3
"""
XYZ Beauty — data-driven analysis & deck builder
Dependencies: pandas, numpy, matplotlib, python-pptx
Run: python xyz_beauty_analysis.py
"""

import os
import sys
from datetime import datetime
from textwrap import dedent

import numpy as np
import pandas as pd
import matplotlib.pyplot as plt

# PowerPoint
from pptx import Presentation
from pptx.util import Inches

# ---------- CONFIG ----------
CSV_PATH = r"C:\Users\HP\beauty_website_dataset.csv" # change if needed
OUTPUT_DIR = r"C:\Users\HP\xyz_output"
DECK_FILENAME = "XYZ_Beauty_Insights.pptx"
# Brands we want to focus on (case brief)
FOCUS_BRANDS = ["Kiehls", "Yves Saint Laurent", "Armani Beauty"]
# Create output dir
os.makedirs(OUTPUT_DIR, exist_ok=True)

# ---------- UTIL HELPERS ----------
def save_fig(fig, name):
    path = os.path.join(OUTPUT_DIR, name)
    fig.savefig(path, bbox_inches="tight")
    plt.close(fig)
    return path

def safe_to_numeric(series):
    return pd.to_numeric(series, errors="coerce")

def normalize_colnames(df):
    df = df.copy()
    df.columns = [c.strip().replace(" ", "_").replace("-", "_").lower() for c in df.columns]
    return df

def canonicalize_brand_column(df, brand_col="brand"):
    """
    Try to map variants of brand names to canonical keys for focus brands:
    returns a column 'brand_canonical' that contains canonical keys where match found,
    otherwise NaN.
    """
    brand_map_candidates = {
        "kiehls": ["kiehls", "kiehl's", "kiehl s", "kiehlʼs"],
        "yves saint laurent": ["yves saint laurent", "ysl", "ysl beauty", "yves-saint-laurent"],
        "armani beauty": ["armani beauty", "giorgio armani beauty", "armani", "giorgio armani"],
    }
    norm = df[brand_col].astype(str).str.lower().str.strip().fillna("")
    df["brand_norm"] = norm
    df["brand_canonical"] = np.nan
    for canon, variants in brand_map_candidates.items():
        mask = False
        for v in variants:
            mask = mask | norm.str.contains(v, na=False)
        df.loc[mask, "brand_canonical"] = canon
    return df

# ---------- LOAD & CLEAN ----------
def load_and_clean(csv_path):
    if not os.path.exists(csv_path):
        raise FileNotFoundError(f"CSV not found at {csv_path}")
    df = pd.read_csv(csv_path)
    df = normalize_colnames(df)

    # Common alternative columns -> standardize names
    rename_map = {}
    if "rating" in df.columns and "ratings" not in df.columns:
        rename_map["rating"] = "ratings"
    if "reviews" in df.columns and "number_of_reviews" not in df.columns:
        rename_map["reviews"] = "number_of_reviews"
    if "mrp" in df.columns and "price" not in df.columns:
        rename_map["mrp"] = "price"
    if "discount_percentage" in df.columns and "discount" not in df.columns:
        rename_map["discount_percentage"] = "discount"
    if rename_map:
        df = df.rename(columns=rename_map)

    # Numeric coercion
    for col in ["ratings", "number_of_reviews", "price", "discount"]:
        if col in df.columns:
            df[col] = safe_to_numeric(df[col])

    # Trim strings in key columns
    for col in ["brand", "category", "name"]:
        if col in df.columns:
            df[col] = df[col].astype(str).str.strip()

    # Drop rows that are fully empty across key columns (name/brand/category/price)
    key_cols = [c for c in ["name", "brand", "category", "price"] if c in df.columns]
    if key_cols:
        df = df.dropna(subset=key_cols, how="all").copy()

    # Normalize brand canonical
    if "brand" in df.columns:
        df = canonicalize_brand_column(df, brand_col="brand")

    return df

# ---------- EDA & KPIs ----------
def compute_summaries(df):
    summary = {}
    summary["n_products"] = len(df)
    summary["columns"] = list(df.columns)
    numeric_cols = [c for c in ["ratings", "number_of_reviews", "price", "discount"] if c in df.columns]
    summary["describe"] = df[numeric_cols].describe().T.round(3) if numeric_cols else pd.DataFrame()

    # Top categories / brands by counts and by reviews
    if "category" in df.columns:
        summary["top_categories_by_count"] = df["category"].value_counts().head(20)
    else:
        summary["top_categories_by_count"] = pd.Series(dtype=int)

    if "brand" in df.columns:
        summary["top_brands_by_count"] = df["brand"].value_counts().head(30)
    else:
        summary["top_brands_by_count"] = pd.Series(dtype=int)

    if set(["category", "number_of_reviews"]).issubset(df.columns):
        summary["top_categories_by_reviews"] = (
            df.groupby("category")["number_of_reviews"].sum().sort_values(ascending=False)
        )
    else:
        summary["top_categories_by_reviews"] = pd.Series(dtype=float)

    if set(["brand", "number_of_reviews"]).issubset(df.columns):
        summary["top_brands_by_reviews"] = (
            df.groupby("brand")["number_of_reviews"].sum().sort_values(ascending=False)
        )
    else:
        summary["top_brands_by_reviews"] = pd.Series(dtype=float)

    # Correlations (Pearson)
    corrs = {}
    cols_for_corr = set(["ratings", "number_of_reviews", "price", "discount"])
    if cols_for_corr.issubset(set(df.columns)):
        corrs["reviews_vs_ratings"] = df["number_of_reviews"].corr(df["ratings"])
        corrs["reviews_vs_price"] = df["number_of_reviews"].corr(df["price"])
        corrs["reviews_vs_discount"] = df["number_of_reviews"].corr(df["discount"])
        corrs["ratings_vs_price"] = df["ratings"].corr(df["price"])
        corrs["ratings_vs_discount"] = df["ratings"].corr(df["discount"])
    summary["correlations"] = corrs

    return summary

# ---------- VISUALIZATIONS ----------
def create_charts(df):
    chart_files = []

    # Ratings distribution
    if "ratings" in df.columns:
        fig = plt.figure(figsize=(8,4))
        df["ratings"].dropna().plot(kind="hist", bins=20)
        plt.title("Distribution of Ratings")
        plt.xlabel("Ratings")
        plt.ylabel("Frequency")
        chart_files.append(save_fig(fig, "01_ratings_distribution.png"))

    # Price distribution
    if "price" in df.columns:
        fig = plt.figure(figsize=(8,4))
        df["price"].dropna().plot(kind="hist", bins=30)
        plt.title("Distribution of Price")
        plt.xlabel("Price")
        plt.ylabel("Frequency")
        chart_files.append(save_fig(fig, "02_price_distribution.png"))

        # price log histogram if skewed
        try:
            p_log = np.log1p(df["price"].dropna())
            fig = plt.figure(figsize=(8,4))
            plt.hist(p_log, bins=30)
            plt.title("Log(1+Price) Distribution")
            plt.xlabel("log(1+price)")
            plt.ylabel("Frequency")
            chart_files.append(save_fig(fig, "02b_price_log_distribution.png"))
        except Exception:
            pass

    # Discount distribution
    if "discount" in df.columns:
        fig = plt.figure(figsize=(8,4))
        df["discount"].dropna().plot(kind="hist", bins=30)
        plt.title("Distribution of Discount")
        plt.xlabel("Discount (%)")
        plt.ylabel("Frequency")
        chart_files.append(save_fig(fig, "03_discount_distribution.png"))

    # Price vs Ratings scatter
    if set(["price", "ratings"]).issubset(df.columns):
        fig = plt.figure(figsize=(6,4))
        plt.scatter(df["price"], df["ratings"], alpha=0.6)
        plt.title("Price vs Ratings")
        plt.xlabel("Price")
        plt.ylabel("Ratings")
        chart_files.append(save_fig(fig, "04_price_vs_ratings.png"))

    # Discount vs Reviews scatter
    if set(["discount", "number_of_reviews"]).issubset(df.columns):
        fig = plt.figure(figsize=(6,4))
        plt.scatter(df["discount"], df["number_of_reviews"], alpha=0.6)
        plt.title("Discount vs Number of Reviews (Demand Proxy)")
        plt.xlabel("Discount (%)")
        plt.ylabel("Number of Reviews")
        chart_files.append(save_fig(fig, "05_discount_vs_reviews.png"))

    # Top categories by total reviews (horizontal bar)
    if "top_categories_by_reviews" in summary and not summary["top_categories_by_reviews"].empty:
        series = summary["top_categories_by_reviews"].head(10).iloc[::-1]
        fig = plt.figure(figsize=(8,4))
        series.plot(kind="barh")
        plt.title("Top Categories by Total Reviews")
        plt.xlabel("Total Reviews")
        plt.ylabel("Category")
        chart_files.append(save_fig(fig, "06_top_categories_reviews.png"))

    # Top brands by total reviews
    if "top_brands_by_reviews" in summary and not summary["top_brands_by_reviews"].empty:
        series = summary["top_brands_by_reviews"].head(15).iloc[::-1]
        fig = plt.figure(figsize=(8,5))
        series.plot(kind="barh")
        plt.title("Top Brands by Total Reviews")
        plt.xlabel("Total Reviews")
        plt.ylabel("Brand")
        chart_files.append(save_fig(fig, "07_top_brands_reviews.png"))

    # Focus brand price vs rating plots
    focus_chart_files = []
    if "brand_canonical" in df.columns and set(["price", "ratings"]).issubset(df.columns):
        for b in ["kiehls", "yves saint laurent", "armani beauty"]:
            sub = df[df["brand_canonical"] == b]
            if len(sub) >= 5:
                fig = plt.figure(figsize=(6,4))
                plt.scatter(sub["price"], sub["ratings"], alpha=0.6)
                plt.title(f"Price vs Ratings — {b.title()}")
                plt.xlabel("Price")
                plt.ylabel("Ratings")
                fname = f"08_price_vs_ratings_{b.replace(' ','_')}.png"
                focus_chart_files.append(save_fig(fig, fname))

    chart_files.extend(focus_chart_files)
    return chart_files

# ---------- BRAND KPIs & CATEGORY MIX ----------
def brand_kpis_and_category_mix(df):
    brand_kpis = None
    category_mix = None
    if set(["brand", "ratings", "number_of_reviews", "price", "discount"]).issubset(df.columns):
        brand_kpis = (
            df.groupby(df["brand_canonical"].fillna(df["brand"]))
              .agg(products=("name", "count"),
                   avg_rating=("ratings", "mean"),
                   total_reviews=("number_of_reviews", "sum"),
                   median_price=("price", "median"),
                   avg_discount=("discount", "mean"))
              .sort_values("total_reviews", ascending=False)
        )
    if set(["brand_canonical", "category", "number_of_reviews", "ratings"]).issubset(df.columns):
        category_mix = (
            df.dropna(subset=["brand_canonical"])
              .groupby(["brand_canonical", "category"])
              .agg(total_reviews=("number_of_reviews", "sum"),
                   avg_rating=("ratings", "mean"),
                   count=("name", "count"))
              .reset_index()
        )
    return brand_kpis, category_mix

# ---------- PPTX BUILD ----------
def build_deck(summary, chart_files, brand_kpis, category_mix, out_path):
    prs = Presentation()
    title_slide_layout = prs.slide_layouts[0]
    title_and_content_layout = prs.slide_layouts[1]
    content_slide_layout = prs.slide_layouts[5] if len(prs.slide_layouts) > 5 else prs.slide_layouts[1]

    # Title
    slide = prs.slides.add_slide(title_slide_layout)
    slide.shapes.title.text = "Company XYZ — Data-Driven Strategy for Beauty Portfolio"
    subtitle = slide.placeholders[1].text_frame
    subtitle.text = f"Insights & Recommendations — generated {datetime.now().strftime('%Y-%m-%d %H:%M')}"

    # Objective slide
    slide = prs.slides.add_slide(title_and_content_layout)
    slide.shapes.title.text = "Objective & Key Questions"
    tf = slide.placeholders[1].text_frame
    tf.text = ("Goal: Extract insights to guide pricing, product development, and marketing for XYZ's focus brands.")
    bullets = [
        "EDA: distributions, outliers, trends across ratings, reviews, price, discount",
        "Customer preferences: popular categories/brands; drivers of demand (reviews)",
        "Product development: attributes linked with higher satisfaction (ratings)",
        "Commercial levers: pricing/discount bands, category focus, hero SKU identification"
    ]
    for b in bullets:
        p = tf.add_paragraph()
        p.text = b
        p.level = 1

    # Key KPIs slide
    slide = prs.slides.add_slide(title_and_content_layout)
    slide.shapes.title.text = "Key KPIs (Overall)"
    tf = slide.placeholders[1].text_frame
    tf.clear()
    tf.text = f"Dataset size: {summary['n_products']} products"
    if not summary["describe"].empty:
        for col in summary["describe"].index:
            mean = summary["describe"].loc[col, "mean"] if "mean" in summary["describe"].columns else None
            median = df[col].median() if col in df.columns else None
            p = tf.add_paragraph()
            p.text = f"{col}: mean={mean:.2f}  median={median:.2f}" if mean is not None else col
            p.level = 1

    # Correlation slide
    if summary.get("correlations"):
        slide = prs.slides.add_slide(title_and_content_layout)
        slide.shapes.title.text = "Demand & Satisfaction Correlates"
        tf = slide.placeholders[1].text_frame
        tf.text = "Correlation snapshot (Pearson):"
        for k, v in summary["correlations"].items():
            p = tf.add_paragraph()
            p.text = f"{k.replace('_',' ').title()}: {v:.2f}"
            p.level = 1

    # Insert charts as full-width images
    def add_image_slide(title, img_path):
        slide = prs.slides.add_slide(content_slide_layout)
        slide.shapes.title.text = title
        left = Inches(0.5)
        top = Inches(1.3)
        height = Inches(5.0)
        slide.shapes.add_picture(img_path, left, top, height=height)

    title_map = {
        "01_ratings_distribution.png": "Distribution of Ratings",
        "02_price_distribution.png": "Distribution of Price",
        "02b_price_log_distribution.png": "Log(1+Price) Distribution",
        "03_discount_distribution.png": "Distribution of Discount",
        "04_price_vs_ratings.png": "Price vs Ratings",
        "05_discount_vs_reviews.png": "Discount vs Reviews",
        "06_top_categories_reviews.png": "Top Categories by Total Reviews",
        "07_top_brands_reviews.png": "Top Brands by Total Reviews",
    }

    for f in chart_files:
        fname = os.path.basename(f)
        title = title_map.get(fname, fname)
        add_image_slide(title, f)

    # Focus brand KPIs slide
    if brand_kpis is not None and not brand_kpis.empty:
        slide = prs.slides.add_slide(title_and_content_layout)
        slide.shapes.title.text = "Focus Brands — Snapshot"
        tf = slide.placeholders[1].text_frame
        tf.text = "Kiehls, Yves Saint Laurent, Armani Beauty"
        # Filter for our canonical keys
        for idx, row in brand_kpis.iterrows():
            idx_str = str(idx)
            if idx_str.lower() in ["kiehls", "yves saint laurent", "armani beauty"]:
                p = tf.add_paragraph()
                p.text = (f"{idx_str.title()}: products={int(row['products'])}, "
                          f"avg_rating={row['avg_rating']:.2f}, median_price={row['median_price']:.0f}, "
                          f"avg_discount={row['avg_discount']:.1f}%, total_reviews={int(row['total_reviews'])}")
                p.level = 1

    # Category plays
    if category_mix is not None and not category_mix.empty:
        bullets = []
        for b in ["kiehls", "yves saint laurent", "armani beauty"]:
            sub = category_mix[category_mix["brand_canonical"] == b]
            if sub.empty:
                continue
            top_by_reviews = sub.sort_values("total_reviews", ascending=False).head(1).iloc[0]
            top_by_rating = sub.sort_values("avg_rating", ascending=False).head(1).iloc[0]
            bullets.append(f"{b.title()}: top volume category={top_by_reviews['category']} (reviews={int(top_by_reviews['total_reviews'])}); top satisfaction category={top_by_rating['category']} (rating={top_by_rating['avg_rating']:.2f})")
        if bullets:
            slide = prs.slides.add_slide(title_and_content_layout)
            slide.shapes.title.text = "Category Plays for Focus Brands"
            tf = slide.placeholders[1].text_frame
            tf.text = bullets[0]
            for b in bullets[1:]:
                p = tf.add_paragraph()
                p.text = b
                p.level = 1

    # Recommendations slide
    slide = prs.slides.add_slide(title_and_content_layout)
    slide.shapes.title.text = "Recommendations"
    tf = slide.placeholders[1].text_frame
    tf.text = "Commercial & Product Levers"
    recs = [
        "Hero SKUs: double down on top categories by reviews; protect price where ratings are high",
        "Smart discounting: use discounts tactically where reviews are price-elastic; avoid blanket markdowns",
        "Portfolio gaps: launch/extend in categories where competitors over-index but ratings remain high",
        "Review velocity: invest in sampling/UGC to accelerate reviews for high-rated, low-review SKUs",
        "Premium laddering: maintain median price leadership only where ratings sustain; bundle to raise AOV",
    ]
    for r in recs:
        p = tf.add_paragraph()
        p.text = r
        p.level = 1

    # Save deck
    prs.save(out_path)
    return out_path

# ---------- RUN PIPELINE ----------
if __name__ == "__main__":
    print("Loading data...")
    df = load_and_clean(CSV_PATH)
    print(f"Loaded {len(df)} rows and {len(df.columns)} columns")

    print("Computing summaries...")
    summary = compute_summaries(df)

    # Export summary tables as CSVs for reference
    summary_dir = os.path.join(OUTPUT_DIR, "tables")
    os.makedirs(summary_dir, exist_ok=True)
    if isinstance(summary.get("describe"), pd.DataFrame) and not summary["describe"].empty:
        summary["describe"].to_csv(os.path.join(summary_dir, "descriptive_stats.csv"))

    if not summary["top_categories_by_count"].empty:
        summary["top_categories_by_count"].to_frame("product_count").to_csv(os.path.join(summary_dir, "top_categories_by_count.csv"))

    if not summary["top_brands_by_count"].empty:
        summary["top_brands_by_count"].to_frame("product_count").to_csv(os.path.join(summary_dir, "top_brands_by_count.csv"))

    if not summary["top_categories_by_reviews"].empty:
        summary["top_categories_by_reviews"].to_frame("total_reviews").to_csv(os.path.join(summary_dir, "top_categories_by_reviews.csv"))

    if not summary["top_brands_by_reviews"].empty:
        summary["top_brands_by_reviews"].to_frame("total_reviews").to_csv(os.path.join(summary_dir, "top_brands_by_reviews.csv"))

    # Brand KPIs
    print("Computing brand KPIs and category mix...")
    brand_kpis, category_mix = brand_kpis_and_category_mix(df)
    if brand_kpis is not None:
        brand_kpis.to_csv(os.path.join(summary_dir, "brand_kpis.csv"))
    if category_mix is not None:
        category_mix.to_csv(os.path.join(summary_dir, "category_mix.csv"))

    # Create charts
    print("Creating charts...")
    chart_files = create_charts(df)
    print(f"Created {len(chart_files)} chart files in {OUTPUT_DIR}")

    # Build deck
    deck_path = os.path.join(OUTPUT_DIR, DECK_FILENAME)
    print("Building PowerPoint deck...")
    built_deck = build_deck(summary, chart_files, brand_kpis, category_mix, deck_path)
    print(f"Deck saved to: {built_deck}")

    print(dedent(f"""
    Done.
    Outputs (in {OUTPUT_DIR}):
      - Charts: PNG files (prefix 01_..)
      - Deck: {DECK_FILENAME}
      - Tables CSVs: {summary_dir}
    """))

